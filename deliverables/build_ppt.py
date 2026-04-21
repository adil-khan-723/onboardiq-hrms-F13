"""
OnboardIQ Presentation — matches website design system exactly.

Website palette (from index.css):
  --bg-base:       #F7F6F3  (warm off-white — slide background)
  --bg-surface:    #FFFFFF  (cards)
  --bg-muted:      #EFEDE8  (subtle fills)
  --bg-subtle:     #E8E5DE  (dividers)
  --text-primary:  #1C1A17  (headings)
  --text-secondary:#5C5852  (body)
  --text-muted:    #8C8880  (captions)
  --sage:          #4A7C6B  (primary accent)
  --sage-light:    #EBF3EF  (sage chip bg)
  --sage-dark:     #2E5A4D  (dark sage)
  --amber:         #C97B2E  (secondary accent)
  --amber-light:   #FDF3E5
  --coral:         #C25042  (error/warning)
  --blue:          #3D6B9E  (info)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Exact website palette ─────────────────────────────────────────────────────
BG       = RGBColor(0xF7, 0xF6, 0xF3)   # --bg-base
SURFACE  = RGBColor(0xFF, 0xFF, 0xFF)   # --bg-surface
MUTED    = RGBColor(0xEF, 0xED, 0xE8)   # --bg-muted
SUBTLE   = RGBColor(0xE8, 0xE5, 0xDE)   # --bg-subtle
T1       = RGBColor(0x1C, 0x1A, 0x17)   # --text-primary
T2       = RGBColor(0x5C, 0x58, 0x52)   # --text-secondary
TM       = RGBColor(0x8C, 0x88, 0x80)   # --text-muted
SAGE     = RGBColor(0x4A, 0x7C, 0x6B)   # --sage
SAGE_L   = RGBColor(0xEB, 0xF3, 0xEF)   # --sage-light
SAGE_D   = RGBColor(0x2E, 0x5A, 0x4D)   # --sage-dark
AMBER    = RGBColor(0xC9, 0x7B, 0x2E)   # --amber
AMBER_L  = RGBColor(0xFD, 0xF3, 0xE5)   # --amber-light
CORAL    = RGBColor(0xC2, 0x50, 0x42)   # --coral
CORAL_L  = RGBColor(0xFB, 0xF0, 0xEE)   # --coral-light
BLUE     = RGBColor(0x3D, 0x6B, 0x9E)   # --blue
BLUE_L   = RGBColor(0xEB, 0xF1, 0xF8)   # --blue-light

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Primitive helpers ─────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(1)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=Pt(13), bold=False, color=T1,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def multiline(slide, rows, x, y, w, h, size=Pt(12), spacing=Pt(4)):
    """rows = [(text, bold, color)]"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    for i, (t, b, c) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        r = p.add_run(); r.text = t
        r.font.size = size; r.font.bold = b; r.font.color.rgb = c

def bg(slide):
    rect(slide, 0, 0, W, H, fill=BG)

def divider(slide, y, color=SUBTLE):
    rect(slide, Inches(0.55), y, Inches(12.2), Pt(1), fill=color)

def chip(slide, text, x, y, fill=SAGE_L, color=SAGE_D):
    w = Inches(0.12) * len(text) + Inches(0.35)
    rect(slide, x, y, w, Inches(0.28), fill=fill)
    txt(slide, text, x, y, w, Inches(0.28),
        size=Pt(9.5), bold=True, color=color, align=PP_ALIGN.CENTER)
    return x + w + Inches(0.12)

def logo_mark(slide, x, y, size=Inches(0.3)):
    r = slide.shapes.add_shape(1, x, y, size, size)
    r.fill.solid(); r.fill.fore_color.rgb = SAGE
    r.line.fill.background()
    inner_off = size * 0.22
    inner_sz  = size * 0.56
    r2 = slide.shapes.add_shape(1,
        x + inner_off, y + inner_off, inner_sz, inner_sz)
    r2.fill.solid(); r2.fill.fore_color.rgb = SAGE_L
    r2.line.fill.background()

def page_header(slide, title, subtitle=None, tag=None, tag_color=SAGE, tag_fill=SAGE_L):
    logo_mark(slide, Inches(0.55), Inches(0.32))
    txt(slide, "OnboardIQ", Inches(0.95), Inches(0.32), Inches(3), Inches(0.3),
        size=Pt(11), bold=True, color=T1)
    if tag:
        rect(slide, Inches(11.5), Inches(0.27), Inches(1.28), Inches(0.28), fill=tag_fill)
        txt(slide, tag, Inches(11.5), Inches(0.27), Inches(1.28), Inches(0.28),
            size=Pt(9), bold=True, color=tag_color, align=PP_ALIGN.CENTER)
    divider(slide, Inches(0.72))
    txt(slide, title, Inches(0.55), Inches(0.85), Inches(10), Inches(0.65),
        size=Pt(28), bold=True, color=T1)
    if subtitle:
        txt(slide, subtitle, Inches(0.55), Inches(1.45), Inches(11), Inches(0.35),
            size=Pt(13), color=T2)
    divider(slide, Inches(1.82) if subtitle else Inches(1.6))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)

# Sage accent strip left
rect(s, 0, 0, Inches(0.04), H, fill=SAGE)

# Subtle background block right side
rect(s, Inches(8.8), Inches(0), Inches(4.53), H, fill=MUTED)

# Logo mark — large version
logo_mark(s, Inches(0.65), Inches(1.5), Inches(0.55))

# WordMark
txt(s, "OnboardIQ", Inches(1.32), Inches(1.5), Inches(5), Inches(0.55),
    size=Pt(14), bold=True, color=T1)

# Main headline
txt(s, "Smart Employee\nOnboarding &\nIdentity Service",
    Inches(0.65), Inches(2.2), Inches(8), Inches(2.8),
    size=Pt(48), bold=True, color=T1)

# Underline accent
rect(s, Inches(0.65), Inches(4.82), Inches(1.1), Inches(0.04), fill=SAGE)

# Subtitle
txt(s, "From offer letter acceptance to Day 1 login — fully automated on AWS.",
    Inches(0.65), Inches(5.05), Inches(7.8), Inches(0.45),
    size=Pt(14), color=T2)

# Tech chips
px = Inches(0.65)
for label, f, c in [
    ("React 18", SAGE_L, SAGE_D),
    ("AWS Lambda", AMBER_L, AMBER),
    ("Step Functions", BLUE_L, BLUE),
    ("DynamoDB", MUTED, T2),
    ("Cognito", CORAL_L, CORAL),
]:
    px = chip(s, label, px, Inches(5.85), fill=f, color=c)

# Right side info
txt(s, "Internship Project", Inches(9.2), Inches(2.0), Inches(3.8), Inches(0.35),
    size=Pt(10), color=TM, italic=True)
txt(s, "AWS Serverless\nFull Stack Application", Inches(9.2), Inches(2.5), Inches(3.8),
    Inches(0.7), size=Pt(22), bold=True, color=T1)

for i, item in enumerate([
    "9 Lambda Functions",
    "4 DynamoDB Tables",
    "Step Functions Workflow",
    "React Frontend on S3",
]):
    y = Inches(3.45) + i * Inches(0.48)
    rect(s, Inches(9.2), y + Inches(0.09), Inches(0.06), Inches(0.06), fill=SAGE)
    txt(s, item, Inches(9.38), y, Inches(3.5), Inches(0.38), size=Pt(12.5), color=T2)

txt(s, "Region: ap-south-1 (Mumbai)",
    Inches(9.2), Inches(5.55), Inches(3.8), Inches(0.28),
    size=Pt(10), color=TM)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, Inches(0.04), H, fill=AMBER)
page_header(s, "The Problem We're Solving",
            "Traditional employee onboarding is manual, fragmented, and slow.",
            tag="PROBLEM", tag_color=AMBER, tag_fill=AMBER_L)

cols = [
    ("Before OnboardIQ", CORAL, CORAL_L, [
        "HR manually emails every new hire",
        "Documents collected over WhatsApp or email",
        "No visibility into who completed what",
        "IT provisions accounts days late",
        "Managers forget to greet new hires",
        "No audit trail or compliance record",
    ]),
    ("With OnboardIQ", SAGE_D, SAGE_L, [
        "Automated emails triggered by workflow engine",
        "Secure encrypted S3 upload portal",
        "Real-time HR dashboard with stage tracking",
        "Cognito account created automatically on submit",
        "Manager introduction email sent on joining date",
        "Every action timestamped in DynamoDB",
    ]),
]

for ci, (heading, hclr, hfill, items) in enumerate(cols):
    x = Inches(0.55) + ci * Inches(6.3)
    rect(s, x, Inches(2.0), Inches(5.9), Inches(0.38), fill=hfill)
    txt(s, heading, x + Inches(0.18), Inches(2.04), Inches(5.5), Inches(0.32),
        size=Pt(11), bold=True, color=hclr)
    for i, item in enumerate(items):
        y = Inches(2.55) + i * Inches(0.72)
        rect(s, x, y, Inches(5.9), Inches(0.62), fill=SURFACE,
             line=SUBTLE, line_w=Pt(0.75))
        dot_c = CORAL if ci == 0 else SAGE
        rect(s, x + Inches(0.18), y + Inches(0.24), Inches(0.07), Inches(0.07), fill=dot_c)
        txt(s, item, x + Inches(0.36), y + Inches(0.12), Inches(5.4), Inches(0.42),
            size=Pt(12), color=T2)

divider(s, Inches(6.87))
txt(s, "OnboardIQ replaces all manual steps with a single AWS Step Functions workflow.",
    Inches(0.55), Inches(7.0), Inches(12.2), Inches(0.32),
    size=Pt(11.5), color=TM, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT IT DOES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, Inches(0.04), H, fill=SAGE)
page_header(s, "What OnboardIQ Does",
            "Two portals. One automated pipeline. Zero manual steps.",
            tag="OVERVIEW")

# New Hire column
txt(s, "New Hire Portal", Inches(0.55), Inches(2.0), Inches(6), Inches(0.38),
    size=Pt(14), bold=True, color=SAGE_D)

steps = [
    ("1", "Fill in your profile",   "Name · Role · Department · Joining date"),
    ("2", "Upload 3 documents",     "ID Proof · Degree Certificate · Offer Letter"),
    ("3", "Sign 5 policies",        "Progress bar fills as each policy is acknowledged"),
    ("4", "Manager intro note",     "Optional message to your reporting manager"),
    ("5", "Done",                   "Credentials arrive in your email inbox"),
]
for i, (num, title, sub) in enumerate(steps):
    y = Inches(2.52) + i * Inches(0.82)
    rect(s, Inches(0.55), y, Inches(0.36), Inches(0.36), fill=SAGE)
    txt(s, num, Inches(0.55), y, Inches(0.36), Inches(0.36),
        size=Pt(11), bold=True, color=SURFACE, align=PP_ALIGN.CENTER)
    txt(s, title, Inches(1.05), y, Inches(5.2), Inches(0.25),
        size=Pt(12.5), bold=True, color=T1)
    txt(s, sub, Inches(1.05), y + Inches(0.26), Inches(5.2), Inches(0.22),
        size=Pt(11), color=T2)

# Divider between columns
rect(s, Inches(6.55), Inches(1.98), Pt(1), Inches(5.3), fill=SUBTLE)

# AWS column
txt(s, "AWS Does Automatically", Inches(6.8), Inches(2.0), Inches(6.1), Inches(0.38),
    size=Pt(14), bold=True, color=T1)

aws_items = [
    ("Employee record created in DynamoDB", SAGE),
    ("Step Functions execution starts",     SAGE),
    ("S3 validates each document upload",   SAGE),
    ("Cognito account provisioned",         SAGE),
    ("Credentials emailed via SES",         AMBER),
    ("Policy sign-off timestamped",         SAGE),
    ("Manager intro email sent",            SAGE),
    ("Employee status → ACTIVE",            SAGE_D),
]
for i, (item, dot) in enumerate(aws_items):
    y = Inches(2.52) + i * Inches(0.54)
    rect(s, Inches(6.8), y + Inches(0.12), Inches(0.07), Inches(0.07), fill=dot)
    txt(s, item, Inches(7.02), y, Inches(5.9), Inches(0.36), size=Pt(12.5), color=T2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, Inches(0.04), H, fill=SAGE)
page_header(s, "System Architecture",
            "Fully serverless — no servers to manage, maintain, or scale.",
            tag="ARCHITECTURE")

layers = [
    ("Frontend",      "React 18 + Vite · Hosted on S3 Static Website",          SAGE,  SAGE_L,  Inches(2.0)),
    ("API Gateway",   "REST API · CORS enabled · Routes to Lambda",               BLUE,  BLUE_L,  Inches(2.9)),
    ("Lambda",        "9 Functions · Node.js 20 · ARM64 · 256MB",                AMBER, AMBER_L, Inches(3.8)),
    ("Step Functions","4-stage state machine: Doc → IT → Policy → Manager Intro", T2,    MUTED,   Inches(4.7)),
]

for name, desc, lclr, lfill, y in layers:
    rect(s, Inches(0.55), y, Inches(12.2), Inches(0.72), fill=lfill,
         line=SUBTLE, line_w=Pt(0.75))
    rect(s, Inches(0.55), y, Inches(0.04), Inches(0.72), fill=lclr)
    txt(s, name, Inches(0.8), y + Inches(0.08), Inches(2.2), Inches(0.28),
        size=Pt(11.5), bold=True, color=lclr)
    txt(s, desc, Inches(0.8), y + Inches(0.36), Inches(11.5), Inches(0.28),
        size=Pt(11), color=T2)

# Connector arrows
for ay in [Inches(2.72), Inches(3.62), Inches(4.52)]:
    txt(s, "↓", Inches(6.4), ay, Inches(0.4), Inches(0.28),
        size=Pt(11), color=TM, align=PP_ALIGN.CENTER)

# Data layer
rect(s, Inches(0.55), Inches(5.55), Inches(12.2), Inches(0.06), fill=SUBTLE)
txt(s, "DATA LAYER", Inches(0.55), Inches(5.68), Inches(2), Inches(0.26),
    size=Pt(9), bold=True, color=TM)

data = [
    ("DynamoDB\n4 Tables",        SAGE,  SAGE_L),
    ("S3 Bucket\nEncrypted",      AMBER, AMBER_L),
    ("Cognito\nUser Pool",        BLUE,  BLUE_L),
    ("SES\nEmail",                SAGE,  MUTED),
    ("SNS\nAlerts",               AMBER, AMBER_L),
]
for di, (label, dc, df) in enumerate(data):
    dx = Inches(0.55) + di * Inches(2.46)
    rect(s, dx, Inches(5.98), Inches(2.28), Inches(0.98), fill=df,
         line=SUBTLE, line_w=Pt(0.75))
    rect(s, dx, Inches(5.98), Inches(0.04), Inches(0.98), fill=dc)
    txt(s, label, dx + Inches(0.18), Inches(6.08), Inches(2.0), Inches(0.78),
        size=Pt(11.5), bold=True, color=dc)

txt(s, "Infrastructure as Code: AWS SAM + CloudFormation · Stack: hrms-onboarding-stack · ap-south-1 (Mumbai)",
    Inches(0.55), Inches(7.1), Inches(12.2), Inches(0.26),
    size=Pt(9.5), color=TM, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — AWS SERVICES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, Inches(0.04), H, fill=AMBER)
page_header(s, "AWS Services Used",
            "9 managed services — zero servers to provision or maintain.",
            tag="SERVICES", tag_color=AMBER, tag_fill=AMBER_L)

services = [
    ("Lambda",        "9 functions · Node.js 20 · ARM64",             AMBER, AMBER_L),
    ("Step Functions","4-stage orchestration state machine",           T2,    MUTED),
    ("DynamoDB",      "4 tables · 3 GSIs · pay-per-request",          BLUE,  BLUE_L),
    ("S3",            "Encrypted document storage + static hosting",   SAGE,  SAGE_L),
    ("Cognito",       "User pool · auto-provisions new hire accounts", CORAL, CORAL_L),
    ("SES",           "Transactional emails — welcome, reminders",     AMBER, AMBER_L),
    ("SNS",           "HR team alerts when all documents received",    SAGE,  SAGE_L),
    ("API Gateway",   "REST API routing all frontend requests",        BLUE,  BLUE_L),
    ("SAM / CF",      "Full infra as code · one-command deploy",       T2,    MUTED),
]

cols_per_row = 3
for i, (name, desc, sc, sf) in enumerate(services):
    col = i % cols_per_row
    row = i // cols_per_row
    x = Inches(0.55) + col * Inches(4.15)
    y = Inches(2.05) + row * Inches(1.6)
    rect(s, x, y, Inches(3.95), Inches(1.42), fill=SURFACE,
         line=SUBTLE, line_w=Pt(0.75))
    rect(s, x, y, Inches(0.04), Inches(1.42), fill=sc)
    rect(s, x, y, Inches(3.95), Inches(0.38), fill=sf)
    txt(s, name, x + Inches(0.18), y + Inches(0.07), Inches(3.6), Inches(0.28),
        size=Pt(12), bold=True, color=sc)
    txt(s, desc, x + Inches(0.18), y + Inches(0.52), Inches(3.6), Inches(0.72),
        size=Pt(11.5), color=T2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ER DIAGRAM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, Inches(0.04), H, fill=SAGE)
page_header(s, "Database Schema",
            "Amazon DynamoDB · 4 tables · 3 Global Secondary Indexes",
            tag="ER DIAGRAM")

tables = [
    {
        "name": "hrms-employees",
        "x": Inches(0.55), "y": Inches(2.05), "w": Inches(3.1), "h": Inches(4.85),
        "color": SAGE, "fill": SAGE_L,
        "pk": "employee_id  PK",
        "fields": [
            ("first_name, last_name, full_name", False),
            ("email", False),
            ("phone, department, role", False),
            ("employment_type, joining_date", False),
            ("manager_id  →  FK", True),
            ("cognito_user_id", False),
            ("status  (ONBOARDING / ACTIVE)", False),
            ("created_at, updated_at", False),
        ]
    },
    {
        "name": "hrms-onboarding-workflows",
        "x": Inches(4.0), "y": Inches(2.05), "w": Inches(3.1), "h": Inches(3.2),
        "color": AMBER, "fill": AMBER_L,
        "pk": "workflow_id  PK",
        "fields": [
            ("employee_id  GSI →  FK", True),
            ("execution_arn", False),
            ("current_stage", False),
            ("overall_status", False),
            ("reminder_attempt", False),
            ("started_at, updated_at", False),
        ]
    },
    {
        "name": "hrms-onboarding-stages",
        "x": Inches(4.0), "y": Inches(5.5), "w": Inches(3.1), "h": Inches(1.78),
        "color": BLUE, "fill": BLUE_L,
        "pk": "stage_id  PK",
        "fields": [
            ("workflow_id  GSI →  FK", True),
            ("stage_name, status", False),
            ("signed_at, completed_at", False),
        ]
    },
    {
        "name": "hrms-documents",
        "x": Inches(7.45), "y": Inches(2.05), "w": Inches(3.1), "h": Inches(4.4),
        "color": CORAL, "fill": CORAL_L,
        "pk": "document_id  PK",
        "fields": [
            ("employee_id  GSI →  FK", True),
            ("doc_type", False),
            ("status  (PENDING / UPLOADED)", False),
            ("s3_key, content_type", False),
            ("verified  (boolean)", False),
            ("validation_errors", False),
            ("uploaded_at, updated_at", False),
        ]
    },
]

for t in tables:
    c = t["color"]; f = t["fill"]
    x, y, w, h = t["x"], t["y"], t["w"], t["h"]
    # Card
    rect(s, x, y, w, h, fill=SURFACE, line=SUBTLE, line_w=Pt(0.75))
    # Header strip
    rect(s, x, y, w, Inches(0.38), fill=f)
    rect(s, x, y, Inches(0.04), h, fill=c)
    txt(s, t["name"], x + Inches(0.14), y + Inches(0.06), w - Inches(0.2), Inches(0.26),
        size=Pt(10), bold=True, color=c)
    # PK row
    rect(s, x, y + Inches(0.38), w, Inches(0.3), fill=MUTED)
    txt(s, "🔑  " + t["pk"], x + Inches(0.14), y + Inches(0.4),
        w - Inches(0.2), Inches(0.26), size=Pt(9.5), bold=True, color=AMBER)
    # Fields
    for fi, (field, is_fk) in enumerate(t["fields"]):
        fy = y + Inches(0.38) + Inches(0.3) + fi * Inches(0.3)
        fc = BLUE if is_fk else T2
        txt(s, field, x + Inches(0.14), fy, w - Inches(0.2), Inches(0.28),
            size=Pt(9.5), color=fc)

# Relationship labels
rel_labels = [
    (Inches(3.68), Inches(3.15), "1 → 1"),
    (Inches(7.14), Inches(3.15), "1 → many"),
    (Inches(4.85), Inches(5.22), "1 → 4"),
]
for rx, ry, label in rel_labels:
    rect(s, rx, ry, Inches(0.72), Inches(0.24), fill=MUTED, line=SUBTLE, line_w=Pt(0.5))
    txt(s, label, rx, ry, Inches(0.72), Inches(0.24),
        size=Pt(8.5), bold=True, color=T2, align=PP_ALIGN.CENTER)

txt(s, "↺ manager_id self-joins back to employee_id for org hierarchy",
    Inches(0.55), Inches(7.1), Inches(5), Inches(0.26),
    size=Pt(9.5), color=TM, italic=True)

txt(s, "Paste the Mermaid diagram in deliverables/02-er-diagram.md at mermaid.live for a rendered version",
    Inches(5.0), Inches(7.1), Inches(7.7), Inches(0.26),
    size=Pt(9.5), color=TM, italic=True, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — STEP FUNCTIONS WORKFLOW
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, Inches(0.04), H, fill=SAGE)
page_header(s, "The Onboarding Workflow Engine",
            "AWS Step Functions — event-driven, auditable, no servers polling.",
            tag="WORKFLOW")

stages = [
    ("Stage 1", "Document\nCollection",
     ["Checks DynamoDB for", "all 3 uploads", "Retries up to 3×"],
     SAGE, SAGE_L),
    ("Stage 2", "IT\nProvisioning",
     ["Creates Cognito user", "Emails temp credentials", "via SES"],
     AMBER, AMBER_L),
    ("Stage 3", "Policy\nSign-off",
     ["Checks signed_at", "timestamp in DynamoDB", "for 5 policies"],
     BLUE, BLUE_L),
    ("Stage 4", "Manager\nIntro",
     ["Emails new hire", "Emails manager", "Status → ACTIVE"],
     CORAL, CORAL_L),
]

box_x = [Inches(0.55), Inches(3.65), Inches(6.75), Inches(9.85)]
for i, (stage, title, details, sc, sf) in enumerate(stages):
    x = box_x[i]
    rect(s, x, Inches(2.05), Inches(2.88), Inches(4.1), fill=SURFACE,
         line=SUBTLE, line_w=Pt(0.75))
    rect(s, x, Inches(2.05), Inches(2.88), Inches(0.34), fill=sf)
    rect(s, x, Inches(2.05), Inches(0.04), Inches(4.1), fill=sc)
    txt(s, stage, x + Inches(0.16), Inches(2.1), Inches(2.6), Inches(0.24),
        size=Pt(9.5), bold=True, color=sc)
    txt(s, title, x + Inches(0.16), Inches(2.45), Inches(2.6), Inches(0.7),
        size=Pt(18), bold=True, color=T1)
    for di, detail in enumerate(details):
        dy = Inches(3.28) + di * Inches(0.4)
        rect(s, x + Inches(0.16), dy + Inches(0.1),
             Inches(0.06), Inches(0.06), fill=sc)
        txt(s, detail, x + Inches(0.32), dy, Inches(2.4), Inches(0.35),
            size=Pt(11.5), color=T2)
    # Complete badge
    rect(s, x + Inches(0.5), Inches(5.78), Inches(1.88), Inches(0.26), fill=sf)
    txt(s, "✓  COMPLETE", x + Inches(0.5), Inches(5.78), Inches(1.88), Inches(0.26),
        size=Pt(9.5), bold=True, color=sc, align=PP_ALIGN.CENTER)

# Arrows
for ax in [Inches(3.44), Inches(6.54), Inches(9.64)]:
    txt(s, "→", ax, Inches(3.72), Inches(0.3), Inches(0.4),
        size=Pt(16), color=TM, align=PP_ALIGN.CENTER)

divider(s, Inches(6.25))
txt(s, "Wait states pause the execution every 30s (demo) / 24h (production) — zero compute cost while waiting.",
    Inches(0.55), Inches(6.38), Inches(12.2), Inches(0.28),
    size=Pt(11), color=TM, italic=True, align=PP_ALIGN.CENTER)

txt(s, "If a stage is incomplete after 3 reminders → marked Overdue → workflow fails gracefully.",
    Inches(0.55), Inches(6.75), Inches(12.2), Inches(0.28),
    size=Pt(11), color=TM, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — LIVE EXECUTION PROOF
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, Inches(0.04), H, fill=SAGE)
page_header(s, "Live Execution Proof",
            "Real AWS Step Functions execution — Arjun Sharma · EMP-4FEA92A3 · 2026-04-08",
            tag="EXECUTION")

# Meta row
rect(s, Inches(0.55), Inches(1.98), Inches(12.2), Inches(0.44), fill=SAGE_L,
     line=SUBTLE, line_w=Pt(0.75))
metas = [
    ("Status", "✓  SUCCEEDED"),
    ("Duration", "66 seconds"),
    ("Stages", "4 / 4 complete"),
    ("Time", "15:11 – 15:12 IST"),
]
for mi, (label, val) in enumerate(metas):
    mx = Inches(0.85) + mi * Inches(3.0)
    txt(s, label, mx, Inches(2.02), Inches(1.2), Inches(0.2),
        size=Pt(8.5), color=TM, bold=True)
    txt(s, val, mx, Inches(2.2), Inches(2.5), Inches(0.2),
        size=Pt(10.5), color=SAGE_D, bold=True)

# Timeline events
events = [
    ("15:11:53", "DocumentCollection entered",    "Stage 1 starts",                        SAGE,  False),
    ("15:11:54", "WaitForDocuments",              "30-second wait — documents uploading",   TM,    True),
    ("15:12:25", "DocumentCollection  ✓",         "All 3 documents received",               SAGE,  False),
    ("15:12:25", "ITProvisioning entered",        "Stage 2 — Cognito user creation",        AMBER, False),
    ("15:12:27", "ITProvisioning  ✓",             "Cognito account created",                AMBER, False),
    ("15:12:27", "PolicySignoff entered",         "Stage 3 — checks signed_at",             BLUE,  False),
    ("15:12:27", "WaitForPolicySignoff",          "30-second wait",                         TM,    True),
    ("15:12:58", "PolicySignoff  ✓",              "All 5 policies signed",                  BLUE,  False),
    ("15:12:58", "ManagerIntro entered",          "Stage 4 — intro emails dispatched",      CORAL, False),
    ("15:12:59", "OnboardingComplete  ✓",         "EXECUTION SUCCEEDED — status → ACTIVE",  SAGE_D,False),
]

for i, (ts, state, desc, ec, dim) in enumerate(events):
    y = Inches(2.62) + i * Inches(0.46)
    if i % 2 == 0:
        rect(s, Inches(0.55), y, Inches(12.2), Inches(0.46), fill=MUTED)
    clr = TM if dim else T2
    txt(s, ts,    Inches(0.72), y + Inches(0.1), Inches(1.1), Inches(0.3),
        size=Pt(10), color=TM)
    rect(s, Inches(1.88), y + Inches(0.18), Inches(0.06), Inches(0.06),
         fill=TM if dim else ec)
    txt(s, state, Inches(2.05), y + Inches(0.1), Inches(3.5), Inches(0.3),
        size=Pt(10.5), bold=not dim, color=TM if dim else ec)
    txt(s, desc,  Inches(5.65), y + Inches(0.1), Inches(6.9), Inches(0.3),
        size=Pt(10.5), color=clr)

divider(s, Inches(7.12))
txt(s, "Open AWS Console → Step Functions → hrms-onboarding-workflow → Executions to see the visual graph",
    Inches(0.55), Inches(7.24), Inches(12.2), Inches(0.22),
    size=Pt(9.5), color=TM, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — DEMO WALKTHROUGH A (NEW HIRE)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, Inches(0.04), H, fill=SAGE)
page_header(s, "Demo Walkthrough — Part A",
            "New Hire Experience · follow these steps exactly during the presentation.",
            tag="DEMO  A")

rect(s, Inches(0.55), Inches(1.98), Inches(12.2), Inches(0.34), fill=SAGE_L,
     line=SUBTLE, line_w=Pt(0.75))
txt(s, "App URL:  hrms-onboarding-frontend-dev.s3-website.ap-south-1.amazonaws.com",
    Inches(0.72), Inches(2.02), Inches(12.0), Inches(0.26),
    size=Pt(10.5), bold=False, color=SAGE_D)

steps_a = [
    ("A1", "Open the App",
     "Navigate to the URL above → splash animation plays → landing page shows two cards.",
     SAGE, SAGE_L),
    ("A2", "Enter New Hire Portal",
     "Click 'New Hire Portal' → another splash → Step 1 of 5 appears.",
     SAGE, SAGE_L),
    ("A3", "Fill Profile",
     "First: Priya  ·  Last: Menon  ·  Email: priya.menon@example.com  ·  Dept: Engineering  ·  Role: Frontend Engineer  →  Continue",
     SAGE, SAGE_L),
    ("A4", "Upload Documents",
     "Click each of the 3 cards → upload any PDF → green Uploaded badge appears on each → Continue",
     AMBER, AMBER_L),
    ("A5", "Sign 5 Policies",
     "Check all 5 boxes one by one — watch the progress bar fill to 100% → Continue",
     AMBER, AMBER_L),
    ("A6", "Manager Intro",
     "Type a note (optional): 'What does the first week look like?' → Submit & Complete ✓",
     BLUE, BLUE_L),
    ("A7", "Completion Screen",
     "Completion screen shows all stage badges. Click '← Back to Home'.",
     SAGE, SAGE_L),
]

for i, (tag, title, detail, tc, tf) in enumerate(steps_a):
    y = Inches(2.48) + i * Inches(0.68)
    rect(s, Inches(0.55), y, Inches(0.36), Inches(0.36), fill=tf, line=tc, line_w=Pt(1))
    txt(s, tag, Inches(0.55), y, Inches(0.36), Inches(0.36),
        size=Pt(9), bold=True, color=tc, align=PP_ALIGN.CENTER)
    txt(s, title, Inches(1.05), y, Inches(2.8), Inches(0.26),
        size=Pt(11.5), bold=True, color=T1)
    txt(s, detail, Inches(1.05), y + Inches(0.26), Inches(11.6), Inches(0.3),
        size=Pt(10.5), color=T2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DEMO WALKTHROUGH B (HR DASHBOARD)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, Inches(0.04), H, fill=AMBER)
page_header(s, "Demo Walkthrough — Part B",
            "HR Admin Dashboard · PIN-gated · real-time pipeline view.",
            tag="DEMO  B", tag_color=AMBER, tag_fill=AMBER_L)

rect(s, Inches(0.55), Inches(1.98), Inches(12.2), Inches(0.34), fill=AMBER_L,
     line=SUBTLE, line_w=Pt(0.75))
txt(s, "HR PIN:  1234   (auto-unlocks on the 4th digit — no Enter needed · resets every navigation)",
    Inches(0.72), Inches(2.02), Inches(12.0), Inches(0.26),
    size=Pt(10.5), color=AMBER, bold=True)

steps_b = [
    ("B1", "Open HR Dashboard",
     "From landing page click 'HR Admin Dashboard' → splash → PIN gate appears.",
     AMBER, AMBER_L),
    ("B2", "Enter PIN",
     "Type 1-2-3-4 — on the 4th digit the screen fades out automatically. Dashboard fades in.",
     AMBER, AMBER_L),
    ("B3", "Stats Bar",
     "5 metric cards at top: Total · In Progress · Completed · Docs Pending · Avg Completion %",
     SAGE, SAGE_L),
    ("B4", "Search for Your Hire",
     "Type 'Priya' in the search bar → table filters live to show Priya Menon's row.",
     SAGE, SAGE_L),
    ("B5", "Expand Detail Panel",
     "Click the row → panel slides open: 4 stages with timestamps + 3 document upload badges.",
     SAGE, SAGE_L),
    ("B6", "Refresh",
     "Click Refresh button (top right) → spinner → data reloads from live API.",
     BLUE, BLUE_L),
    ("B7", "Status Filters",
     "Use filter pills: All · In Progress · Completed · Pending — each filters the table.",
     BLUE, BLUE_L),
]

for i, (tag, title, detail, tc, tf) in enumerate(steps_b):
    y = Inches(2.48) + i * Inches(0.68)
    rect(s, Inches(0.55), y, Inches(0.36), Inches(0.36), fill=tf, line=tc, line_w=Pt(1))
    txt(s, tag, Inches(0.55), y, Inches(0.36), Inches(0.36),
        size=Pt(9), bold=True, color=tc, align=PP_ALIGN.CENTER)
    txt(s, title, Inches(1.05), y, Inches(2.8), Inches(0.26),
        size=Pt(11.5), bold=True, color=T1)
    txt(s, detail, Inches(1.05), y + Inches(0.26), Inches(11.6), Inches(0.3),
        size=Pt(10.5), color=T2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — DEMO WALKTHROUGH C (AWS CONSOLE)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, Inches(0.04), H, fill=BLUE)
page_header(s, "Demo Walkthrough — Part C",
            "AWS Console evidence — show evaluators the live backend.",
            tag="DEMO  C", tag_color=BLUE, tag_fill=BLUE_L)

console_steps = [
    ("C1", "Step Functions",
     "AWS Console → Step Functions → hrms-onboarding-workflow → Executions tab\n"
     "Click latest SUCCEEDED → Graph view (all states green) → Events tab (full timeline)",
     SAGE),
    ("C2", "DynamoDB — Employees",
     "DynamoDB → hrms-employees → Explore items → find your employee\n"
     "Show: status = ACTIVE · cognito_user_id populated · all fields correct",
     BLUE),
    ("C3", "DynamoDB — Stages",
     "DynamoDB → hrms-onboarding-stages → filter by workflow_id\n"
     "Show: all 4 rows with status = COMPLETE and completed_at timestamps",
     BLUE),
    ("C4", "S3 — Documents",
     "S3 → hrms-onboarding-documents-736786104206-dev → documents/{employee_id}/\n"
     "Show: 3 uploaded files — ID_PROOF · DEGREE_CERTIFICATE · OFFER_LETTER",
     AMBER),
    ("C5", "S3 — Encryption",
     "S3 bucket → Properties → Default encryption → AES-256 server-side encryption enabled",
     AMBER),
]

for i, (tag, title, detail, tc) in enumerate(console_steps):
    y = Inches(2.05) + i * Inches(1.02)
    rect(s, Inches(0.55), y, Inches(12.2), Inches(0.92), fill=SURFACE,
         line=SUBTLE, line_w=Pt(0.75))
    rect(s, Inches(0.55), y, Inches(0.04), Inches(0.92), fill=tc)
    rect(s, Inches(0.55), y, Inches(12.2), Inches(0.32), fill=MUTED)
    txt(s, tag + "  ·  " + title, Inches(0.75), y + Inches(0.05),
        Inches(11.8), Inches(0.24), size=Pt(11), bold=True, color=tc)
    txt(s, detail, Inches(0.75), y + Inches(0.36), Inches(11.8), Inches(0.52),
        size=Pt(10.5), color=T2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — COST ESTIMATE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, Inches(0.04), H, fill=SAGE)
page_header(s, "Cost Estimate — 50 Hires / Month",
            "AWS ap-south-1 pricing · April 2026 · all major services within free tier at this scale.",
            tag="COST")

costs = [
    ("Lambda",         "750 invocations/mo",      "$0.00", "Covered by 1M free tier",   SAGE,  SAGE_L),
    ("Step Functions", "1,650 transitions/mo",    "$0.00", "First 4,000 free",           T2,    MUTED),
    ("DynamoDB",       "~6,500 reads + writes",   "$0.00", "25GB always-free tier",      BLUE,  BLUE_L),
    ("S3",             "75MB · 150 PUTs",         "$0.004","$0.023/GB storage",          AMBER, AMBER_L),
    ("SES",            "400 emails/mo",           "$0.04", "$0.10 per 1,000 emails",     SAGE,  SAGE_L),
    ("Cognito",        "50 MAUs",                 "$0.00", "First 50,000 MAU free",      CORAL, CORAL_L),
    ("SNS",            "50 notifications",        "$0.00", "First 1M free",              T2,    MUTED),
    ("API Gateway",    "~200 requests/mo",        "$0.001","$3.50 per million requests", BLUE,  BLUE_L),
]

# Header row
rect(s, Inches(0.55), Inches(2.0), Inches(12.2), Inches(0.36), fill=MUTED)
for xi, hdr in zip(
    [Inches(0.72), Inches(3.1), Inches(6.8), Inches(9.0)],
    ["Service", "Monthly Usage", "Cost", "Pricing Basis"]
):
    txt(s, hdr, xi, Inches(2.06), Inches(2.5), Inches(0.26),
        size=Pt(10), bold=True, color=T2)

for i, (name, usage, cost, basis, sc, sf) in enumerate(costs):
    y = Inches(2.36) + i * Inches(0.54)
    if i % 2 == 0:
        rect(s, Inches(0.55), y, Inches(12.2), Inches(0.54), fill=SURFACE)
    rect(s, Inches(0.55), y, Inches(0.04), Inches(0.54), fill=sc)
    txt(s, name,  Inches(0.72), y + Inches(0.12), Inches(2.2), Inches(0.3),
        size=Pt(11.5), bold=True, color=sc)
    txt(s, usage, Inches(3.1),  y + Inches(0.12), Inches(3.5), Inches(0.3),
        size=Pt(11), color=T2)
    txt(s, cost,  Inches(6.8),  y + Inches(0.12), Inches(1.8), Inches(0.3),
        size=Pt(12), bold=True, color=SAGE_D if cost == "$0.00" else T1)
    txt(s, basis, Inches(9.0),  y + Inches(0.12), Inches(3.5), Inches(0.3),
        size=Pt(10.5), color=TM)

# Total
rect(s, Inches(0.55), Inches(6.74), Inches(12.2), Inches(0.5), fill=SAGE_L,
     line=SAGE, line_w=Pt(1))
txt(s, "Total monthly cost:", Inches(0.72), Inches(6.82), Inches(5), Inches(0.32),
    size=Pt(12), bold=True, color=SAGE_D)
txt(s, "~$0.05 / month  (≈ ₹4)",
    Inches(5.5), Inches(6.82), Inches(5), Inches(0.32),
    size=Pt(15), bold=True, color=SAGE_D)
txt(s, "Scales to 5,000 hires/month for ~$5.00",
    Inches(10.5), Inches(6.86), Inches(2.5), Inches(0.26),
    size=Pt(10), color=TM, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, Inches(0.04), H, fill=SAGE)

# Right block
rect(s, Inches(8.5), 0, Inches(4.83), H, fill=MUTED)

logo_mark(s, Inches(0.65), Inches(1.8), Inches(0.6))
txt(s, "OnboardIQ", Inches(1.38), Inches(1.8), Inches(5), Inches(0.6),
    size=Pt(15), bold=True, color=T1)

txt(s, "From offer letter\nto Day 1 login.",
    Inches(0.65), Inches(2.6), Inches(7.5), Inches(1.8),
    size=Pt(48), bold=True, color=T1)

rect(s, Inches(0.65), Inches(4.3), Inches(1.4), Inches(0.04), fill=SAGE)

txt(s, "Fully automated. Fully serverless. Built on AWS.",
    Inches(0.65), Inches(4.55), Inches(7.5), Inches(0.4),
    size=Pt(15), color=T2)

divider(s, Inches(5.3))

links = [
    ("Live App", "hrms-onboarding-frontend-dev.s3-website.ap-south-1.amazonaws.com"),
    ("HR PIN",   "1234"),
    ("Region",   "ap-south-1  ·  Mumbai"),
    ("GitHub",   "github.com/adil-khan-723/onboardiq-hrms-F13"),
]
for li, (label, val) in enumerate(links):
    ly = Inches(5.55) + li * Inches(0.42)
    txt(s, label, Inches(0.65), ly, Inches(1.3), Inches(0.28),
        size=Pt(10), bold=True, color=TM)
    txt(s, val,   Inches(2.05), ly, Inches(5.8), Inches(0.28),
        size=Pt(10.5), color=T1)

txt(s, "Thank You",
    Inches(0.65), Inches(7.0), Inches(7.5), Inches(0.38),
    size=Pt(11), color=TM, italic=True)

# Right block content
txt(s, "Built with", Inches(9.0), Inches(1.8), Inches(3.8), Inches(0.3),
    size=Pt(10), color=TM)
services_r = [
    ("AWS Lambda",        AMBER, AMBER_L),
    ("Step Functions",    T2,    MUTED),
    ("DynamoDB",          BLUE,  BLUE_L),
    ("S3",                SAGE,  SAGE_L),
    ("Cognito",           CORAL, CORAL_L),
    ("SES  ·  SNS",       AMBER, AMBER_L),
    ("API Gateway",       BLUE,  BLUE_L),
    ("React 18 + Vite",   SAGE,  SAGE_L),
    ("AWS SAM",           T2,    MUTED),
]
for si, (sname, sc, sf) in enumerate(services_r):
    sy = Inches(2.25) + si * Inches(0.52)
    rect(s, Inches(9.0), sy, Inches(3.6), Inches(0.38), fill=sf,
         line=SUBTLE, line_w=Pt(0.5))
    rect(s, Inches(9.0), sy, Inches(0.04), Inches(0.38), fill=sc)
    txt(s, sname, Inches(9.18), sy + Inches(0.07), Inches(3.3), Inches(0.26),
        size=Pt(11.5), bold=True, color=sc)


SS_DIR = "/Users/oggy/F13/project 4/docs/observability/screenshots"

def ss_slide(title, subtitle, tag, img_path, caption, accent=SAGE, tag_fill=SAGE_L):
    s = prs.slides.add_slide(blank)
    bg(s)
    rect(s, 0, 0, Inches(0.04), H, fill=accent)
    page_header(s, title, subtitle, tag=tag, tag_color=accent, tag_fill=tag_fill)
    # Screenshot image — full width below header
    pic = s.shapes.add_picture(img_path, Inches(0.55), Inches(1.95),
                               Inches(12.2), Inches(5.1))
    # Caption bar
    rect(s, Inches(0.55), Inches(7.08), Inches(12.2), Inches(0.3), fill=tag_fill)
    txt(s, caption, Inches(0.72), Inches(7.1), Inches(12.0), Inches(0.26),
        size=Pt(9.5), color=accent, italic=True)
    return s

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CLOUDWATCH DASHBOARD (screenshot)
# ══════════════════════════════════════════════════════════════════════════════
ss_slide(
    "CloudWatch Dashboard — Live",
    "hrms-onboarding-observability · 4 widgets · all Lambda functions · ap-south-1",
    "OBSERVABILITY",
    f"{SS_DIR}/dashboard.png",
    "Dashboard: Lambda Invocations · Lambda Errors · p99 Duration · Active Alarms — all 12 alarms green except demo alarm",
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — TRIGGERED ALARM (screenshot)
# ══════════════════════════════════════════════════════════════════════════════
ss_slide(
    "Triggered Alarm Demo",
    "hrms-alarm-simulation-error-rate · ALARM state · 6 deliberate errors fired",
    "ALARM DEMO",
    f"{SS_DIR}/alarm.png",
    "Alarm triggered at 21:30 IST · 6 errors > threshold of 1 · state: OK → ALARM · SNS email dispatched automatically",
    accent=CORAL, tag_fill=CORAL_L,
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — SNS EMAIL RECEIPT (screenshot)
# ══════════════════════════════════════════════════════════════════════════════
ss_slide(
    "SNS Alarm Email — Received",
    "AWS CloudWatch → SNS → adilk81054@gmail.com · state change OK → ALARM",
    "ALERT EMAIL",
    f"{SS_DIR}/email.png",
    "Email received at 16:24:21 UTC · Alarm: hrms-alarm-simulation-error-rate · SNS Topic: hrms-hr-notifications · Account: 736786104206",
    accent=AMBER, tag_fill=AMBER_L,
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — TRUSTED ADVISOR (screenshot)
# ══════════════════════════════════════════════════════════════════════════════
ss_slide(
    "AWS Trusted Advisor",
    "Free checks run — 0 critical findings · all security posture checks passed",
    "SECURITY",
    f"{SS_DIR}/trusted.png",
    "Basic support plan: full check suite requires Business/Enterprise plan · S3 public access blocked · DynamoDB SSE enabled · IAM scoped",
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — IAM HARDENING SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, Inches(0.04), H, fill=CORAL)
page_header(s, "IAM Security Hardening",
            "Before & after — SES permissions scoped across all Lambda functions.",
            tag="IAM AUDIT", tag_color=CORAL, tag_fill=CORAL_L)

# Before column
rect(s, Inches(0.55), Inches(2.0), Inches(5.9), Inches(0.36), fill=CORAL_L)
txt(s, "BEFORE  —  Overly broad", Inches(0.72), Inches(2.05), Inches(5.6), Inches(0.26),
    size=Pt(11), bold=True, color=CORAL)

before_items = [
    ("hrms-stage-document-collection", "ses:SendEmail / ses:SendRawEmail", "Resource: '*'"),
    ("hrms-create-employee",           "ses:SendEmail / ses:SendRawEmail", "Resource: '*'"),
    ("hrms-stage-it-provisioning",     "ses:SendEmail / ses:SendRawEmail", "Resource: '*'"),
    ("hrms-stage-policy-signoff",      "ses:SendEmail / ses:SendRawEmail", "Resource: '*'"),
    ("hrms-stage-manager-intro",       "ses:SendEmail / ses:SendRawEmail", "Resource: '*'"),
    ("hrms-send-reminder",             "ses:SendEmail / ses:SendRawEmail", "Resource: '*'"),
    ("hrms-document-upload-trigger",   "ses:SendEmail / ses:SendRawEmail", "Resource: '*'"),
]
for i, (fn, action, resource) in enumerate(before_items):
    y = Inches(2.44) + i * Inches(0.6)
    rect(s, Inches(0.55), y, Inches(5.9), Inches(0.52), fill=SURFACE,
         line=SUBTLE, line_w=Pt(0.75))
    rect(s, Inches(0.55), y, Inches(0.04), Inches(0.52), fill=CORAL)
    txt(s, fn,       Inches(0.72), y + Inches(0.04), Inches(5.5), Inches(0.22),
        size=Pt(9.5), bold=True, color=T1)
    txt(s, f"{action}  ·  {resource}", Inches(0.72), y + Inches(0.26), Inches(5.5), Inches(0.22),
        size=Pt(9), color=CORAL)

# After column
rect(s, Inches(6.85), Inches(2.0), Inches(5.9), Inches(0.36), fill=SAGE_L)
txt(s, "AFTER  —  Least privilege", Inches(7.02), Inches(2.05), Inches(5.6), Inches(0.26),
    size=Pt(11), bold=True, color=SAGE_D)

after_resource = "arn:aws:ses:<region>:<account>:identity/*"
for i, (fn, action, _) in enumerate(before_items):
    y = Inches(2.44) + i * Inches(0.6)
    rect(s, Inches(6.85), y, Inches(5.9), Inches(0.52), fill=SURFACE,
         line=SUBTLE, line_w=Pt(0.75))
    rect(s, Inches(6.85), y, Inches(0.04), Inches(0.52), fill=SAGE)
    txt(s, fn,       Inches(7.02), y + Inches(0.04), Inches(5.5), Inches(0.22),
        size=Pt(9.5), bold=True, color=T1)
    txt(s, f"{action}  ·  {after_resource}", Inches(7.02), y + Inches(0.26), Inches(5.5), Inches(0.22),
        size=Pt(9), color=SAGE_D)

divider(s, Inches(6.7))
txt(s, "Also added: xray:PutTraceSegments + xray:PutTelemetryRecords on all Lambdas · Cognito + Step Functions were already resource-scoped ✓",
    Inches(0.55), Inches(6.82), Inches(12.2), Inches(0.28),
    size=Pt(10), color=TM, italic=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — COST ESTIMATE (observability)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, Inches(0.04), H, fill=SAGE)
page_header(s, "Cost Estimate — 50 Users / Day",
            "All services itemised · ap-south-1 · April 2026 · includes observability layer.",
            tag="COST", tag_color=SAGE, tag_fill=SAGE_L)

costs2 = [
    ("Lambda",         "31,000 invocations/mo",   "$0.00",  "Within 400K GB-s free tier",      SAGE,  SAGE_L),
    ("API Gateway",    "~4,650 requests/mo",       "$0.02",  "$3.50 per million requests",       BLUE,  BLUE_L),
    ("DynamoDB",       "~231K reads + writes",     "$0.42",  "On-demand · SSE enabled",          BLUE,  BLUE_L),
    ("S3",             "150 uploads · 300MB",      "$0.01",  "$0.025/GB · versioning on",        AMBER, AMBER_L),
    ("SES",            "12,000 emails/mo",         "$1.20",  "$0.10 per 1,000 emails",           SAGE,  SAGE_L),
    ("Cognito",        "50 MAUs",                  "$0.00",  "First 50,000 MAU free",            CORAL, CORAL_L),
    ("Step Functions", "30,000 transitions/mo",    "$0.65",  "$0.000025 per transition",         T2,    MUTED),
    ("CloudWatch",     "Logs + Dashboard + Alarms","$4.45",  "1 dashboard $3 · 12 alarms $1.20", AMBER, AMBER_L),
    ("X-Ray",          "31,000 traces/mo",         "$0.00",  "Within 100K free-tier traces",     SAGE,  SAGE_L),
    ("SNS",            "~200 notifications/mo",    "$0.00",  "First 1M free",                    T2,    MUTED),
]

rect(s, Inches(0.55), Inches(2.0), Inches(12.2), Inches(0.34), fill=MUTED)
for xi, hdr in zip([Inches(0.72), Inches(3.0), Inches(6.6), Inches(8.6)],
                   ["Service", "Monthly Usage", "Cost", "Pricing Basis"]):
    txt(s, hdr, xi, Inches(2.05), Inches(2.5), Inches(0.26),
        size=Pt(10), bold=True, color=T2)

for i, (name, usage, cost, basis, sc, sf) in enumerate(costs2):
    y = Inches(2.34) + i * Inches(0.48)
    if i % 2 == 0:
        rect(s, Inches(0.55), y, Inches(12.2), Inches(0.48), fill=SURFACE)
    rect(s, Inches(0.55), y, Inches(0.04), Inches(0.48), fill=sc)
    txt(s, name,  Inches(0.72), y + Inches(0.1), Inches(2.1), Inches(0.28),
        size=Pt(11), bold=True, color=sc)
    txt(s, usage, Inches(3.0),  y + Inches(0.1), Inches(3.4), Inches(0.28),
        size=Pt(10.5), color=T2)
    txt(s, cost,  Inches(6.6),  y + Inches(0.1), Inches(1.6), Inches(0.28),
        size=Pt(11.5), bold=True, color=SAGE_D if cost == "$0.00" else T1)
    txt(s, basis, Inches(8.6),  y + Inches(0.1), Inches(3.9), Inches(0.28),
        size=Pt(10), color=TM)

rect(s, Inches(0.55), Inches(7.04), Inches(12.2), Inches(0.36), fill=SAGE_L,
     line=SAGE, line_w=Pt(1))
txt(s, "Total monthly cost at 50 users/day:", Inches(0.72), Inches(7.12),
    Inches(5), Inches(0.26), size=Pt(11), bold=True, color=SAGE_D)
txt(s, "~$6.75 / month", Inches(6.0), Inches(7.1), Inches(4), Inches(0.28),
    size=Pt(14), bold=True, color=SAGE_D)
txt(s, "X-Ray + CloudWatch observability included",
    Inches(10.0), Inches(7.14), Inches(2.6), Inches(0.22),
    size=Pt(9.5), color=TM, italic=True)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/oggy/F13/project 4/deliverables/OnboardIQ-Presentation.pptx"
prs.save(out)
print("Saved:", out)
print(f"Slides: {len(prs.slides)}")
